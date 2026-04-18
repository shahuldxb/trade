import os
import time
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional
from urllib.parse import urlparse

from dotenv import load_dotenv
from tqdm import tqdm

from langchain_openai import AzureOpenAIEmbeddings
from langchain_postgres import PGVector
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader, PDFMinerLoader, TextLoader

try:
    import pandas as pd
except Exception:
    pd = None

try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from azure.ai.documentintelligence import DocumentIntelligenceClient
    from azure.core.credentials import AzureKeyCredential
except Exception:
    DocumentIntelligenceClient = None
    AzureKeyCredential = None

try:
    import pikepdf
except Exception:
    pikepdf = None

try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None


logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
logger = logging.getLogger("ingest_pdfs")


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif"}
TEXT_EXTENSIONS = {".txt", ".md"}
CSV_EXTENSIONS = {".csv"}
EXCEL_EXTENSIONS = {".xls", ".xlsx"}
WORD_EXTENSIONS = {".doc", ".docx"}
PPT_EXTENSIONS = {".ppt", ".pptx"}
PDF_EXTENSIONS = {".pdf"}

SUPPORTED_EXTENSIONS = (
    IMAGE_EXTENSIONS
    | TEXT_EXTENSIONS
    | CSV_EXTENSIONS
    | EXCEL_EXTENSIONS
    | WORD_EXTENSIONS
    | PPT_EXTENSIONS
    | PDF_EXTENSIONS
)


# ---------- helpers ----------

def get_env(name: str, required: bool = True, default: Optional[str] = None) -> str:
    value = os.getenv(name, default)
    if required and not value:
        raise ValueError(f"Missing env var: {name}")
    return value


def get_env_int(name: str, default: int) -> int:
    value = os.getenv(name)
    if value is None or value == "":
        return default
    try:
        return int(value)
    except ValueError:
        logger.warning("Invalid int for %s=%s. Using default=%s", name, value, default)
        return default


def get_env_float(name: str, default: float) -> float:
    value = os.getenv(name)
    if value is None or value == "":
        return default
    try:
        return float(value)
    except ValueError:
        logger.warning("Invalid float for %s=%s. Using default=%s", name, value, default)
        return default


def get_env_bool(name: str, default: bool = False) -> bool:
    value = os.getenv(name)
    if value is None or value == "":
        return default
    return value.strip().lower() in {"1", "true", "yes", "y", "on"}


def _try_loader(loader_cls, path: str, **kwargs) -> List[Document]:
    try:
        return loader_cls(path, **kwargs).load()
    except Exception as exc:
        logger.warning("[%s] %s: %s", loader_cls.__name__, os.path.basename(path), exc)
        return []


def _sanitize_pdf(path: str) -> str:
    if pikepdf is None:
        return path
    try:
        os.makedirs(".tmp_sanitized", exist_ok=True)
        out_path = os.path.join(".tmp_sanitized", os.path.basename(path))
        with pikepdf.open(path) as pdf:
            pdf.save(out_path, linearize=True)
        return out_path
    except Exception as exc:
        logger.warning("[pikepdf] could not sanitize %s: %s", os.path.basename(path), exc)
        return path


def _normalize_docs(
    docs: List[Document],
    source_path: str,
    file_type: str,
    ocr: bool,
    loader_name: str
) -> List[Document]:
    for doc in docs:
        if doc.metadata is None:
            doc.metadata = {}
        doc.metadata.setdefault("source", source_path)
        doc.metadata["file_type"] = file_type
        doc.metadata["ocr"] = ocr
        doc.metadata["loader"] = loader_name
    return docs


def _filter_empty(docs: List[Document]) -> List[Document]:
    out = []
    for doc in docs:
        if doc.page_content and doc.page_content.strip():
            out.append(doc)
    return out


def _sanitize_documents(docs: List[Document]) -> List[Document]:
    for doc in docs:
        if doc.page_content:
            doc.page_content = doc.page_content.replace("\x00", "")
        if doc.metadata:
            for key, value in list(doc.metadata.items()):
                if isinstance(value, str):
                    doc.metadata[key] = value.replace("\x00", "")
    return docs


def _sanitize_azure_endpoint(endpoint: str) -> str:
    try:
        parsed = urlparse(endpoint)
        if parsed.scheme and parsed.netloc:
            return f"{parsed.scheme}://{parsed.netloc}"
    except Exception:
        pass
    return endpoint


def iter_supported_files(docs_dir: str) -> List[Path]:
    root = Path(docs_dir)
    files: List[Path] = []
    for path in root.rglob("*"):
        if not path.is_file():
            continue
        if path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(path)
    return sorted(files)


def batched(items: List[Document], batch_size: int) -> Iterable[List[Document]]:
    for idx in range(0, len(items), batch_size):
        yield items[idx:idx + batch_size]


@dataclass
class IngestionConfig:
    docs_dir: str
    db_url: str
    endpoint: str
    api_key: str
    api_version: str
    emb_deployment: str
    collection: str

    pdf_password_env: str = "PDF_PASSWORD"

    ocr_enabled: bool = True
    ocr_min_total_chars: int = 50
    ocr_min_avg_chars_per_page: int = 20
    ocr_min_text_page_ratio: float = 0.2
    ocr_model: str = "prebuilt-layout"
    ocr_pdf_via_images: bool = True

    di_endpoint: Optional[str] = None
    di_key: Optional[str] = None

    chunk_size: int = 1000
    chunk_overlap: int = 200
    batch_size: int = 32


class DocumentExtractor:
    def __init__(self, config: IngestionConfig):
        self.config = config
        self._di_client: Optional[DocumentIntelligenceClient] = None
        self._di_warned = False

        self._handlers: Dict[str, callable] = {
            ".pdf": self._extract_pdf,
            ".txt": self._extract_text,
            ".md": self._extract_text,
            ".csv": self._extract_csv,
            ".xls": self._extract_excel,
            ".xlsx": self._extract_excel,
            ".doc": self._extract_word,
            ".docx": self._extract_word,
            ".ppt": self._extract_ppt,
            ".pptx": self._extract_ppt,
        }
        for ext in IMAGE_EXTENSIONS:
            self._handlers[ext] = self._extract_image

    def extract(self, path: Path) -> List[Document]:
        ext = path.suffix.lower()
        handler = self._handlers.get(ext, self._extract_unknown)
        docs = handler(path)
        docs = _filter_empty(docs)

        if ext in IMAGE_EXTENSIONS:
            return docs

        if self._should_ocr(docs):
            ocr_docs = self._extract_with_document_intelligence(path)
            ocr_docs = _filter_empty(ocr_docs)
            if ocr_docs:
                return ocr_docs

        return docs

    def _should_ocr(self, docs: List[Document]) -> bool:
        if not self.config.ocr_enabled:
            return False
        if not self._di_ready():
            if not self._di_warned:
                logger.warning("Azure Document Intelligence not configured; OCR fallback disabled.")
                self._di_warned = True
            return False
        if not docs:
            return True

        total_chars = sum(len((d.page_content or "").strip()) for d in docs)
        avg_chars = total_chars / max(len(docs), 1)
        text_pages = sum(1 for d in docs if (d.page_content or "").strip())
        text_page_ratio = text_pages / max(len(docs), 1)

        return (
            total_chars < self.config.ocr_min_total_chars
            or avg_chars < self.config.ocr_min_avg_chars_per_page
            or text_page_ratio < self.config.ocr_min_text_page_ratio
        )

    def _extract_pdf(self, path: Path) -> List[Document]:
        pdf_password = os.getenv(self.config.pdf_password_env, "") or None

        docs = _try_loader(PyPDFLoader, str(path), password=pdf_password)
        loader_name = "PyPDFLoader"
        if not docs:
            docs = _try_loader(PDFMinerLoader, str(path))
            loader_name = "PDFMinerLoader"

        if not docs:
            repaired = _sanitize_pdf(str(path))
            if repaired != str(path):
                docs = _try_loader(PyPDFLoader, repaired, password=pdf_password)
                loader_name = "PyPDFLoader"
                if not docs:
                    docs = _try_loader(PDFMinerLoader, repaired)
                    loader_name = "PDFMinerLoader"

        return _normalize_docs(docs, str(path), "pdf", False, loader_name)

    def _extract_text(self, path: Path) -> List[Document]:
        docs = []
        try:
            docs = TextLoader(str(path), encoding="utf-8").load()
            loader = "TextLoader"
        except Exception:
            try:
                content = Path(path).read_text(encoding="utf-8", errors="ignore")
                docs = [Document(page_content=content, metadata={})]
                loader = "TextFile"
            except Exception as exc:
                logger.warning("[Text] %s: %s", os.path.basename(path), exc)
                loader = "TextFile"
        return _normalize_docs(docs, str(path), "text", False, loader)

    def _extract_csv(self, path: Path) -> List[Document]:
        content = ""
        if pd is not None:
            try:
                df = pd.read_csv(path, dtype=str, keep_default_na=False)
                content = df.to_csv(index=False)
                loader = "pandas"
            except Exception as exc:
                logger.warning("[CSV pandas] %s: %s", os.path.basename(path), exc)
                loader = "csv"
        else:
            loader = "csv"

        if not content:
            try:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            except Exception as exc:
                logger.warning("[CSV] %s: %s", os.path.basename(path), exc)
                content = ""

        docs = [Document(page_content=content, metadata={})] if content else []
        return _normalize_docs(docs, str(path), "csv", False, loader)

    def _extract_excel(self, path: Path) -> List[Document]:
        if pd is None:
            logger.warning("[Excel] pandas/openpyxl not available: %s", os.path.basename(path))
            return []
        try:
            sheets = pd.read_excel(path, sheet_name=None, dtype=str)
        except Exception as exc:
            logger.warning("[Excel] %s: %s", os.path.basename(path), exc)
            return []

        docs: List[Document] = []
        for sheet_name, df in sheets.items():
            content = df.to_csv(index=False)
            doc = Document(page_content=content, metadata={"sheet": sheet_name})
            docs.append(doc)

        return _normalize_docs(docs, str(path), "excel", False, "pandas")

    def _extract_word(self, path: Path) -> List[Document]:
        if DocxDocument is None:
            logger.warning("[Word] python-docx not available: %s", os.path.basename(path))
            return []

        try:
            docx = DocxDocument(str(path))
        except Exception as exc:
            logger.warning("[Word] %s: %s", os.path.basename(path), exc)
            return []

        parts: List[str] = []
        for p in docx.paragraphs:
            if p.text and p.text.strip():
                parts.append(p.text)

        for table in docx.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text)
                if row_text:
                    parts.append(row_text)

        content = "\n".join(parts)
        docs = [Document(page_content=content, metadata={})] if content else []
        return _normalize_docs(docs, str(path), "word", False, "python-docx")

    def _extract_ppt(self, path: Path) -> List[Document]:
        if Presentation is None:
            logger.warning("[PPT] python-pptx not available: %s", os.path.basename(path))
            return []

        try:
            pres = Presentation(str(path))
        except Exception as exc:
            logger.warning("[PPT] %s: %s", os.path.basename(path), exc)
            return []

        docs: List[Document] = []
        for idx, slide in enumerate(pres.slides, start=1):
            parts: List[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    parts.append(text)
            content = "\n".join(parts)
            if content:
                docs.append(Document(page_content=content, metadata={"slide": idx}))

        return _normalize_docs(docs, str(path), "ppt", False, "python-pptx")

    def _extract_image(self, path: Path) -> List[Document]:
        return self._extract_with_document_intelligence(path)

    def _extract_unknown(self, path: Path) -> List[Document]:
        logger.warning("Unsupported file type: %s", path)
        return []

    def _di_ready(self) -> bool:
        return bool(
            self.config.di_endpoint
            and self.config.di_key
            and DocumentIntelligenceClient is not None
            and AzureKeyCredential is not None
        )

    def _get_di_client(self) -> Optional[DocumentIntelligenceClient]:
        if not self._di_ready():
            return None
        if self._di_client is None:
            self._di_client = DocumentIntelligenceClient(
                endpoint=self.config.di_endpoint,
                credential=AzureKeyCredential(self.config.di_key)
            )
        return self._di_client

    def _extract_with_document_intelligence(self, path: Path) -> List[Document]:
        client = self._get_di_client()
        if client is None:
            return []

        file_type = path.suffix.lower().lstrip(".") or "unknown"

        if file_type == "pdf" and self.config.ocr_pdf_via_images:
            return self._extract_pdf_with_images(path)

        try:
            with open(path, "rb") as f:
                poller = client.begin_analyze_document(
                    self.config.ocr_model,
                    body=f,
                    content_type="application/octet-stream"
                )
                result = poller.result()
        except Exception as exc:
            logger.warning("[DI OCR] %s: %s", os.path.basename(path), exc)
            return []

        docs: List[Document] = []
        pages = getattr(result, "pages", None) or []
        if pages:
            for idx, page in enumerate(pages, start=1):
                lines = getattr(page, "lines", None) or []
                text = "\n".join(
                    line.content for line in lines if getattr(line, "content", None)
                )
                docs.append(
                    Document(
                        page_content=text,
                        metadata={"page": idx}
                    )
                )
        else:
            content = getattr(result, "content", "") or ""
            if content:
                docs.append(Document(page_content=content, metadata={"page": 1}))

        return _normalize_docs(docs, str(path), file_type, True, "azure-document-intelligence")

    def _extract_pdf_with_images(self, path: Path) -> List[Document]:
        if fitz is None:
            logger.warning("[PyMuPDF] not available, falling back to DI on PDF: %s", os.path.basename(path))
            return self._extract_with_document_intelligence(path)

        client = self._get_di_client()
        if client is None:
            return []

        docs: List[Document] = []
        try:
            pdf_doc = fitz.open(str(path))
        except Exception as exc:
            logger.warning("[PyMuPDF] %s: %s", os.path.basename(path), exc)
            return self._extract_with_document_intelligence(path)

        for page_index in range(len(pdf_doc)):
            try:
                page = pdf_doc[page_index]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                img_bytes = pix.tobytes("jpeg")
            except Exception as exc:
                logger.warning("[PyMuPDF] render failed %s p.%s: %s", os.path.basename(path), page_index + 1, exc)
                continue

            text = self._ocr_image_bytes(client, img_bytes)
            docs.append(Document(page_content=text, metadata={"page": page_index + 1}))

        pdf_doc.close()
        return _normalize_docs(docs, str(path), "pdf", True, "azure-di-image")

    def _ocr_image_bytes(self, client: DocumentIntelligenceClient, image_bytes: bytes) -> str:
        try:
            poller = client.begin_analyze_document(
                self.config.ocr_model,
                body=image_bytes,
                content_type="image/jpeg"
            )
            result = poller.result()
        except Exception as exc:
            logger.warning("[DI OCR] image bytes: %s", exc)
            return ""

        pages = getattr(result, "pages", None) or []
        lines = []
        for page in pages:
            for line in getattr(page, "lines", None) or []:
                if getattr(line, "content", None):
                    lines.append(line.content)
        if lines:
            return "\n".join(lines)
        return (getattr(result, "content", "") or "")


# ---------- main pipeline ----------

def load_config() -> IngestionConfig:
    raw_endpoint = get_env("AZURE_OPENAI_ENDPOINT")
    sanitized_endpoint = _sanitize_azure_endpoint(raw_endpoint)
    if sanitized_endpoint != raw_endpoint:
        logger.warning("AZURE_OPENAI_ENDPOINT contains a path; using base %s", sanitized_endpoint)

    return IngestionConfig(
        docs_dir=get_env("DOCS_DIR"),
        db_url=get_env("PGVECTOR_DATABASE_URL"),
        endpoint=sanitized_endpoint,
        api_key=get_env("AZURE_OPENAI_API_KEY"),
        api_version=get_env("AZURE_OPENAI_API_VERSION"),
        emb_deployment=get_env("AZURE_OPENAI_EMBEDDINGS_DEPLOYMENT"),
        collection=get_env("COLLECTION_NAME"),
        ocr_enabled=get_env_bool("OCR_ENABLED", True),
        ocr_min_total_chars=get_env_int("OCR_MIN_TOTAL_CHARS", 50),
        ocr_min_avg_chars_per_page=get_env_int("OCR_MIN_AVG_CHARS_PER_PAGE", 20),
        ocr_min_text_page_ratio=get_env_float("OCR_MIN_TEXT_PAGE_RATIO", 0.2),
        ocr_model=os.getenv("OCR_MODEL", "prebuilt-layout"),
        ocr_pdf_via_images=get_env_bool("OCR_PDF_VIA_IMAGES", True),
        di_endpoint=os.getenv("AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT"),
        di_key=os.getenv("AZURE_DOCUMENT_INTELLIGENCE_KEY"),
        chunk_size=get_env_int("CHUNK_SIZE", 1000),
        chunk_overlap=get_env_int("CHUNK_OVERLAP", 200),
        batch_size=get_env_int("BATCH_SIZE", 32),
    )


def _build_vectorstore(config: IngestionConfig) -> PGVector:
    embeddings = AzureOpenAIEmbeddings(
        azure_endpoint=config.endpoint,
        api_key=config.api_key,
        api_version=config.api_version,
        deployment=config.emb_deployment,
    )

    return PGVector(
        embeddings=embeddings,
        connection=config.db_url,
        collection_name=config.collection,
        use_jsonb=True,
    )


def ingest_paths(paths: List[Path], config: Optional[IngestionConfig] = None) -> Dict[str, int]:
    config = config or load_config()
    extractor = DocumentExtractor(config)

    raw_docs: List[Document] = []
    for path in paths:
        if not path.exists():
            logger.warning("Path not found: %s", path)
            continue
        docs = extractor.extract(path)
        if docs:
            raw_docs.extend(docs)
        else:
            logger.warning("No content extracted: %s", path)

    raw_docs = _sanitize_documents(_filter_empty(raw_docs))
    if not raw_docs:
        logger.warning("Nothing to ingest for provided paths.")
        return {"docs": 0, "chunks": 0, "inserted": 0}

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.chunk_size,
        chunk_overlap=config.chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", " ", ""],
    )
    chunks = splitter.split_documents(raw_docs)
    for chunk in chunks:
        if chunk.metadata is None:
            chunk.metadata = {}
        chunk.metadata["chunk_size"] = config.chunk_size
        chunk.metadata["chunk_overlap"] = config.chunk_overlap
        chunk.metadata["batch_size"] = config.batch_size
        chunk.metadata["collection_name"] = config.collection

    vectorstore = _build_vectorstore(config)

    inserted = 0
    with tqdm(total=len(chunks), desc="Embedding & inserting", unit="chunk") as pbar:
        for batch in batched(chunks, config.batch_size):
            start_time = time.time()
            ids = vectorstore.add_documents(batch)
            inserted += len(ids)
            elapsed = time.time() - start_time
            pbar.set_postfix({"last_batch_sec": f"{elapsed:.2f}"})
            pbar.update(len(batch))

    return {"docs": len(raw_docs), "chunks": len(chunks), "inserted": inserted}


def ingest_directory(docs_dir: str, config: Optional[IngestionConfig] = None) -> Dict[str, int]:
    files = iter_supported_files(docs_dir)
    logger.info("Found %s supported file(s) in %s", len(files), docs_dir)
    if not files:
        logger.warning("No supported files found in %s", docs_dir)
        return {"docs": 0, "chunks": 0, "inserted": 0}
    return ingest_paths(files, config)


def main() -> None:
    load_dotenv()
    config = load_config()

    docs_dir = Path(config.docs_dir)
    if not docs_dir.exists():
        raise FileNotFoundError(f"DOCS_DIR not found: {docs_dir}")

    logger.info("Loading files from: %s", docs_dir.resolve())
    result = ingest_directory(str(docs_dir), config)
    logger.info(
        "Ingest complete | docs=%s chunks=%s inserted=%s",
        result["docs"],
        result["chunks"],
        result["inserted"],
    )


if __name__ == "__main__":
    main()
